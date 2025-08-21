#!/usr/bin/env python3
"""
Nested Document Processor - Robust extraction and chunking for nested agenda documents
Handles PDFs, Word docs, PowerPoint, and complex layouts with adaptive strategies
"""

import os
import json
import logging
import re
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
from dataclasses import dataclass
import hashlib

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

@dataclass
class NestedDocument:
    """Represents a nested document with extracted content"""
    file_path: str
    file_name: str
    parent_item_id: str
    content_type: str
    extracted_text: str
    chunk_count: int
    extraction_quality: str
    file_size: int
    processing_errors: List[str]

@dataclass
class DocumentChunk:
    """Represents a chunk from a nested document"""
    chunk_id: str
    parent_document: str
    parent_item_id: str
    content: str
    content_type: str
    chunk_type: str
    section_header: Optional[str]
    page_number: Optional[int]
    chunk_index: int
    total_chunks: int

class NestedDocumentProcessor:
    """Robust processor for nested agenda documents"""
    
    def __init__(self):
        self.content_type_patterns = {
            'resolution': [
                r'RESOLUTION.*TOWN OF DAVIE',
                r'BE IT RESOLVED',
                r'WHEREAS',
                r'NOW, THEREFORE'
            ],
            'staff_report': [
                r'staff report',
                r'prepared by:',
                r'executive summary',
                r'background',
                r'recommendation'
            ],
            'exhibit': [
                r'exhibit [a-z]',
                r'map',
                r'plan',
                r'drawing',
                r'plat'
            ],
            'planning_report': [
                r'planning.*board',
                r'zoning.*board',
                r'development.*review',
                r'planning.*report'
            ],
            'development_report': [
                r'development.*review',
                r'broward.*county',
                r'development.*report'
            ]
        }
        
        self.section_headers = {
            'resolution': ['WHEREAS', 'NOW, THEREFORE', 'BE IT RESOLVED', 'PROVIDING FOR'],
            'staff_report': ['Executive Summary', 'Background', 'Analysis', 'Recommendation', 'Fiscal Impact'],
            'exhibit': ['Description', 'Location', 'Details', 'Specifications'],
            'planning_report': ['Project Description', 'Site Analysis', 'Recommendations', 'Conditions'],
            'development_report': ['Review Comments', 'Conditions', 'Recommendations', 'Approval']
        }
    
    def extract_text_robust(self, file_path: str) -> Tuple[str, List[str]]:
        """Extract text from any document format with multiple fallbacks"""
        
        errors = []
        text = ""
        
        try:
            # Try PyMuPDF first (handles 90% of cases)
            import fitz
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                page_text = page.get_text()
                if page_text:
                    text += page_text + "\n"
            doc.close()
            
            if len(text.strip()) > 100:
                return text, errors
                
        except Exception as e:
            errors.append(f"PyMuPDF failed: {str(e)}")
        
        try:
            # Fallback to pdfplumber for complex layouts
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                text = ""
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            
            if len(text.strip()) > 100:
                return text, errors
                
        except Exception as e:
            errors.append(f"pdfplumber failed: {str(e)}")
        
        try:
            # Try python-docx for Word documents
            if file_path.endswith('.docx'):
                import docx
                doc = docx.Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                
                if len(text.strip()) > 100:
                    return text, errors
                    
        except Exception as e:
            errors.append(f"python-docx failed: {str(e)}")
        
        try:
            # Try python-pptx for PowerPoint
            if file_path.endswith('.pptx'):
                import pptx
                prs = pptx.Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                
                if len(text.strip()) > 100:
                    return text, errors
                    
        except Exception as e:
            errors.append(f"python-pptx failed: {str(e)}")
        
        # If all else fails, return minimal text
        errors.append("All extraction methods failed")
        return "Extraction failed - document may be image-based or corrupted", errors
    
    def detect_content_type(self, text: str, file_name: str) -> str:
        """Detect what type of document this is"""
        
        text_lower = text.lower()
        file_lower = file_name.lower()
        
        # Check each content type pattern
        for content_type, patterns in self.content_type_patterns.items():
            for pattern in patterns:
                if re.search(pattern, text_lower, re.IGNORECASE):
                    return content_type
        
        # Fallback based on filename
        if 'resolution' in file_lower:
            return 'resolution'
        elif 'report' in file_lower:
            return 'staff_report'
        elif 'exhibit' in file_lower:
            return 'exhibit'
        elif 'planning' in file_lower:
            return 'planning_report'
        elif 'development' in file_lower:
            return 'development_report'
        
        return 'general_document'
    
    def validate_extraction(self, text: str) -> Tuple[bool, str]:
        """Ensure we got meaningful content"""
        
        if len(text.strip()) < 50:
            return False, "Too short"
        
        # Check for common extraction failures
        if text.count(" ") < 10:
            return False, "No spaces - likely image"
        
        if "error" in text.lower() or "failed" in text.lower():
            return False, "Extraction error"
        
        # Check for meaningful content
        meaningful_words = ['the', 'and', 'or', 'in', 'on', 'at', 'to', 'for', 'of', 'with']
        word_count = len([word for word in text.lower().split() if word in meaningful_words])
        
        if word_count < 5:
            return False, "Insufficient meaningful content"
        
        return True, "Valid"
    
    def chunk_by_sections(self, text: str, section_markers: List[str]) -> List[DocumentChunk]:
        """Chunk by section markers (for resolutions, reports)"""
        
        chunks = []
        lines = text.split('\n')
        current_chunk = []
        current_section = "Introduction"
        chunk_index = 0
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check if this line starts a new section
            is_section_header = any(marker.lower() in line.lower() for marker in section_markers)
            
            if is_section_header and current_chunk:
                # Save current chunk
                chunk_text = '\n'.join(current_chunk)
                if len(chunk_text.strip()) > 50:
                    chunks.append(DocumentChunk(
                        chunk_id=f"nested_chunk_{chunk_index}",
                        parent_document="",  # Will be set later
                        parent_item_id="",   # Will be set later
                        content=chunk_text,
                        content_type="",     # Will be set later
                        chunk_type="section",
                        section_header=current_section,
                        page_number=None,
                        chunk_index=chunk_index,
                        total_chunks=0  # Will be set later
                    ))
                    chunk_index += 1
                
                current_chunk = [line]
                current_section = line
            else:
                current_chunk.append(line)
        
        # Add final chunk
        if current_chunk:
            chunk_text = '\n'.join(current_chunk)
            if len(chunk_text.strip()) > 50:
                chunks.append(DocumentChunk(
                    chunk_id=f"nested_chunk_{chunk_index}",
                    parent_document="",
                    parent_item_id="",
                    content=chunk_text,
                    content_type="",
                    chunk_type="section",
                    section_header=current_section,
                    page_number=None,
                    chunk_index=chunk_index,
                    total_chunks=0
                ))
        
        # Set total_chunks for all chunks
        for chunk in chunks:
            chunk.total_chunks = len(chunks)
        
        return chunks
    
    def chunk_by_paragraphs(self, text: str, max_length: int = 1000) -> List[DocumentChunk]:
        """Chunk by paragraphs with size limits"""
        
        chunks = []
        paragraphs = text.split('\n\n')
        current_chunk = []
        current_length = 0
        chunk_index = 0
        
        for paragraph in paragraphs:
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            
            if current_length + len(paragraph) > max_length and current_chunk:
                # Save current chunk
                chunk_text = '\n\n'.join(current_chunk)
                chunks.append(DocumentChunk(
                    chunk_id=f"nested_chunk_{chunk_index}",
                    parent_document="",
                    parent_item_id="",
                    content=chunk_text,
                    content_type="",
                    chunk_type="paragraph",
                    section_header=None,
                    page_number=None,
                    chunk_index=chunk_index,
                    total_chunks=0
                ))
                chunk_index += 1
                current_chunk = [paragraph]
                current_length = len(paragraph)
            else:
                current_chunk.append(paragraph)
                current_length += len(paragraph)
        
        # Add final chunk
        if current_chunk:
            chunk_text = '\n\n'.join(current_chunk)
            chunks.append(DocumentChunk(
                chunk_id=f"nested_chunk_{chunk_index}",
                parent_document="",
                parent_item_id="",
                content=chunk_text,
                content_type="",
                chunk_type="paragraph",
                section_header=None,
                page_number=None,
                chunk_index=chunk_index,
                total_chunks=0
            ))
        
        # Set total_chunks for all chunks
        for chunk in chunks:
            chunk.total_chunks = len(chunks)
        
        return chunks
    
    def chunk_smart(self, text: str, max_length: int = 800, overlap: int = 100) -> List[DocumentChunk]:
        """Smart chunking with overlap for general documents"""
        
        chunks = []
        chunk_index = 0
        
        # Split into sentences first
        sentences = re.split(r'[.!?]+', text)
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            if current_length + len(sentence) > max_length and current_chunk:
                # Save current chunk
                chunk_text = '. '.join(current_chunk) + '.'
                chunks.append(DocumentChunk(
                    chunk_id=f"nested_chunk_{chunk_index}",
                    parent_document="",
                    parent_item_id="",
                    content=chunk_text,
                    content_type="",
                    chunk_type="smart",
                    section_header=None,
                    page_number=None,
                    chunk_index=chunk_index,
                    total_chunks=0
                ))
                chunk_index += 1
                
                # Keep last few sentences for overlap
                overlap_sentences = current_chunk[-3:] if len(current_chunk) > 3 else current_chunk
                current_chunk = overlap_sentences + [sentence]
                current_length = sum(len(s) for s in current_chunk)
            else:
                current_chunk.append(sentence)
                current_length += len(sentence)
        
        # Add final chunk
        if current_chunk:
            chunk_text = '. '.join(current_chunk) + '.'
            chunks.append(DocumentChunk(
                chunk_id=f"nested_chunk_{chunk_index}",
                parent_document="",
                parent_item_id="",
                content=chunk_text,
                content_type="",
                chunk_type="smart",
                section_header=None,
                page_number=None,
                chunk_index=chunk_index,
                total_chunks=0
            ))
        
        # Set total_chunks for all chunks
        for chunk in chunks:
            chunk.total_chunks = len(chunks)
        
        return chunks
    
    def chunk_nested_document(self, text: str, content_type: str) -> List[DocumentChunk]:
        """Chunk based on document type"""
        
        if content_type == "resolution":
            # Legal documents - chunk by sections
            section_markers = self.section_headers.get('resolution', ['WHEREAS', 'NOW, THEREFORE'])
            return self.chunk_by_sections(text, section_markers)
        
        elif content_type == "staff_report":
            # Reports - chunk by headers
            section_markers = self.section_headers.get('staff_report', ['Executive Summary', 'Background'])
            return self.chunk_by_sections(text, section_markers)
        
        elif content_type == "exhibit":
            # Maps/plans - chunk by paragraphs
            return self.chunk_by_paragraphs(text, max_length=1000)
        
        elif content_type in ["planning_report", "development_report"]:
            # Planning reports - chunk by sections
            section_markers = self.section_headers.get(content_type, ['Project Description', 'Recommendations'])
            return self.chunk_by_sections(text, section_markers)
        
        else:
            # General - smart chunking
            return self.chunk_smart(text, max_length=800, overlap=100)
    
    def extract_parent_item_id(self, file_name: str) -> str:
        """Extract parent agenda item ID from filename"""
        # Example: item_8433_583_nested_1.pdf -> 8433
        match = re.search(r'item_(\d+)_\d+_nested', file_name)
        if match:
            return match.group(1)
        return ""
    
    def process_nested_documents(self, agendas_dir: str = "downloads/agendas") -> List[NestedDocument]:
        """Process all nested documents in the agendas directory"""
        
        nested_documents = []
        agendas_path = Path(agendas_dir)
        
        if not agendas_path.exists():
            logger.error(f"Agendas directory not found: {agendas_dir}")
            return nested_documents
        
        # Find all nested files
        nested_files = []
        for meeting_dir in agendas_path.iterdir():
            if meeting_dir.is_dir():
                for file_path in meeting_dir.iterdir():
                    if file_path.is_file() and 'nested' in file_path.name:
                        nested_files.append(file_path)
        
        logger.info(f"Found {len(nested_files)} nested files to process")
        
        for file_path in nested_files:
            try:
                logger.info(f"Processing: {file_path.name}")
                
                # Extract text
                text, errors = self.extract_text_robust(str(file_path))
                
                # Validate extraction
                is_valid, validation_msg = self.validate_extraction(text)
                if not is_valid:
                    errors.append(f"Validation failed: {validation_msg}")
                
                # Detect content type
                content_type = self.detect_content_type(text, file_path.name)
                
                # Extract parent item ID
                parent_item_id = self.extract_parent_item_id(file_path.name)
                
                # Create nested document
                nested_doc = NestedDocument(
                    file_path=str(file_path),
                    file_name=file_path.name,
                    parent_item_id=parent_item_id,
                    content_type=content_type,
                    extracted_text=text,
                    chunk_count=0,  # Will be set after chunking
                    extraction_quality="good" if is_valid else "poor",
                    file_size=file_path.stat().st_size,
                    processing_errors=errors
                )
                
                nested_documents.append(nested_doc)
                
            except Exception as e:
                logger.error(f"Error processing {file_path.name}: {str(e)}")
                continue
        
        logger.info(f"Successfully processed {len(nested_documents)} nested documents")
        return nested_documents
    
    def create_chunks_from_documents(self, nested_documents: List[NestedDocument]) -> List[DocumentChunk]:
        """Create chunks from all nested documents"""
        
        all_chunks = []
        
        for doc in nested_documents:
            try:
                # Chunk the document
                chunks = self.chunk_nested_document(doc.extracted_text, doc.content_type)
                
                # Update chunk metadata
                for chunk in chunks:
                    chunk.parent_document = doc.file_name
                    chunk.parent_item_id = doc.parent_item_id
                    chunk.content_type = doc.content_type
                
                all_chunks.extend(chunks)
                doc.chunk_count = len(chunks)
                
                logger.info(f"Created {len(chunks)} chunks from {doc.file_name}")
                
            except Exception as e:
                logger.error(f"Error chunking {doc.file_name}: {str(e)}")
                continue
        
        logger.info(f"Total chunks created: {len(all_chunks)}")
        return all_chunks
    
    def save_results(self, nested_documents: List[NestedDocument], chunks: List[DocumentChunk], 
                    output_file: str = "nested_documents_results.json"):
        """Save processing results"""
        
        results = {
            "nested_documents": [
                {
                    "file_path": doc.file_path,
                    "file_name": doc.file_name,
                    "parent_item_id": doc.parent_item_id,
                    "content_type": doc.content_type,
                    "extraction_quality": doc.extraction_quality,
                    "file_size": doc.file_size,
                    "chunk_count": doc.chunk_count,
                    "processing_errors": doc.processing_errors
                }
                for doc in nested_documents
            ],
            "chunks": [
                {
                    "chunk_id": chunk.chunk_id,
                    "parent_document": chunk.parent_document,
                    "parent_item_id": chunk.parent_item_id,
                    "content": chunk.content,
                    "content_type": chunk.content_type,
                    "chunk_type": chunk.chunk_type,
                    "section_header": chunk.section_header,
                    "chunk_index": chunk.chunk_index,
                    "total_chunks": chunk.total_chunks
                }
                for chunk in chunks
            ],
            "summary": {
                "total_documents": len(nested_documents),
                "total_chunks": len(chunks),
                "content_types": {},
                "quality_stats": {}
            }
        }
        
        # Add summary statistics
        content_types = {}
        quality_stats = {"good": 0, "poor": 0}
        
        for doc in nested_documents:
            content_types[doc.content_type] = content_types.get(doc.content_type, 0) + 1
            quality_stats[doc.extraction_quality] = quality_stats.get(doc.extraction_quality, 0) + 1
        
        results["summary"]["content_types"] = content_types
        results["summary"]["quality_stats"] = quality_stats
        
        with open(output_file, 'w') as f:
            json.dump(results, f, indent=2, ensure_ascii=False)
        
        logger.info(f"Results saved to {output_file}")
        return results

def main():
    """Main processing function"""
    
    processor = NestedDocumentProcessor()
    
    # Process all nested documents
    logger.info("🚀 Starting nested document processing...")
    nested_documents = processor.process_nested_documents()
    
    # Create chunks
    logger.info("📄 Creating chunks from documents...")
    chunks = processor.create_chunks_from_documents(nested_documents)
    
    # Save results
    logger.info("💾 Saving results...")
    results = processor.save_results(nested_documents, chunks)
    
    # Display summary
    logger.info("=" * 60)
    logger.info("📊 PROCESSING SUMMARY")
    logger.info("=" * 60)
    logger.info(f"Documents processed: {results['summary']['total_documents']}")
    logger.info(f"Chunks created: {results['summary']['total_chunks']}")
    logger.info(f"Quality: {results['summary']['quality_stats']}")
    logger.info(f"Content types: {results['summary']['content_types']}")
    
    return results

if __name__ == "__main__":
    main() 